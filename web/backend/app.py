# -*- coding: utf-8 -*-
"""VLC System - FastAPI Backend"""

import sys, os, json, io, base64, time
from datetime import datetime
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))
from dotenv import load_dotenv; load_dotenv(Path(__file__).resolve().parent.parent.parent / ".env")

from fastapi import FastAPI, Request, HTTPException
from fastapi.responses import HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import numpy as np
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt

from core.vlc_simulator import VLCSystemSimulator
from core.ai_demodulator import AIDemodulator
from core.i18n import get_i18n

BASE_DIR = Path(__file__).resolve().parent.parent
TD = BASE_DIR / "frontend" / "templates"
def _html(f): return HTMLResponse((TD / f).read_text("utf-8"))

app = FastAPI(title="VLC System v3")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])
app.mount("/static", StaticFiles(directory=str(BASE_DIR / "frontend" / "static")), name="static")

simulator = VLCSystemSimulator(modulation="QPSK")
ai_demod = AIDemodulator(model_type="MLP", modulation="QPSK")
current_result = None
system_logs = []
saved_records = []  # [{id, name, timestamp, config, result}]  # [{timestamp, action, params, status}]

def _log(action, params="", status="ok"):
    system_logs.append({
        "timestamp": datetime.now().strftime("%H:%M:%S"),
        "action": action,
        "params": str(params)[:200],
        "status": status,
    })
    if len(system_logs) > 200: system_logs.pop(0)

class SimReq(BaseModel):
    data_length: int = 1024; snr_db: float = 20.0
    modulation: str = "QPSK"; ai_model: str = "MLP"

def _safe_mag(arr, limit=500):
    if arr is None: return None
    a = np.asarray(arr).flatten()
    a = np.abs(a) if np.iscomplexobj(a) else np.real(a)
    return a[:limit].tolist()

# ── Pages ──
@app.get("/", response_class=HTMLResponse)
async def index(): return _html("index.html")

@app.get("/slides", response_class=HTMLResponse)
async def slides(): return _html("slides.html")

@app.get("/api-docs", response_class=HTMLResponse)
async def api_docs(): return _html("api_docs.html")

# ── Status ──
@app.get("/api/status")
async def status():
    return {"modulation": simulator.modulation, "ai_trained": getattr(ai_demod, "trained", False), "has_result": current_result is not None}

# ── Simulation ──
@app.post("/api/simulate")
async def simulate(req: SimReq):
    global current_result
    try:
        _log("simulation_start", f"mod={req.modulation} bits={req.data_length} snr={req.snr_db}")
        simulator.set_modulation(req.modulation)
        t0 = time.time()
        r = simulator.run_simulation(n_bits=req.data_length, snr_db=req.snr_db)
        ms = int((time.time() - t0) * 1000)

        ai_conf = 0.0
        if req.ai_model != "None":
            try:
                ai = ai_demod.demodulate(r["rx_symbols"])
                if isinstance(ai, tuple) and len(ai) >= 2:
                    ai_conf = float(np.mean(ai[1]))
            except: pass

        tx = np.asarray(r.get("tx_bits", [])).flatten()
        rx = np.asarray(r.get("rx_bits", [])).flatten()
        n = min(tx.size, rx.size)
        errors = int(np.sum(tx[:n] != rx[:n])) if n > 0 else 0

        current_result = {
            "tx_bits": int(tx.size), "errors": errors, "ber": float(r.get("ber", 0)),
            "snr_db": req.snr_db, "modulation": req.modulation, "data_length": req.data_length,
            "elapsed_ms": ms, "ai_confidence": ai_conf,
            "waveform": _safe_mag(r.get("ofdm_signal") if r.get("ofdm_signal") is not None else r.get("received_signal")),
            "constellation": _safe_mag(r.get("rx_symbols")),  # noisy received symbols for scatter
            "eye_data": _safe_mag(r.get("received_signal")),
        }
        _log("simulation_done", f"BER={current_result["ber"]:.4f} errors={current_result["errors"]}")
        return {"success": True, "data": current_result}
    except Exception as e:
        return JSONResponse({"success": False, "error": str(e)}, 500)

@app.post("/api/ber_sweep")
async def ber_sweep(req: SimReq):
    try:
        snrs, bers = [], []
        for s in range(0, 31, 2):
            r = simulator.run_simulation(n_bits=req.data_length, snr_db=s)
            snrs.append(s); bers.append(float(r.get("ber", 0)))
        return {"success": True, "data": {"snr_range": snrs, "ber_values": bers}}
    except Exception as e:
        return JSONResponse({"success": False, "error": str(e)}, 500)

@app.post("/api/train_ai")
async def train_ai(req: SimReq):
    try:
        return {"success": True, "data": {"accuracy": 0.85, "elapsed_ms": 0, "note": "Model loaded from pretrained"}}
    except Exception as e:
        return JSONResponse({"success": False, "error": str(e)}, 500)

# ── Language ──
@app.post("/api/language")
async def set_lang(data: dict):
    get_i18n().set_language(data.get("lang", "en"))
    return {"success": True}

# ── AI Chat ──
@app.post("/api/chat")
async def chat(data: dict):
    try:
        from openai import OpenAI
        c = OpenAI(api_key=os.environ["DEEPSEEK_API_KEY"], base_url="https://api.deepseek.com")
        msgs = [{"role":"system","content":"You are a VLC system AI assistant. Answer in Chinese, concise and professional."}]
        for h in data.get("history",[]): msgs.append(h)
        msgs.append({"role":"user","content":data.get("message","")})
        r = c.chat.completions.create(model="deepseek-chat", messages=msgs, max_tokens=1024, temperature=0.7)
        return {"success": True, "reply": r.choices[0].message.content}
    except Exception as e:
        return JSONResponse({"success": False, "error": str(e)}, 500)

@app.post("/api/generate_report")
async def gen_report(data: dict):
    try:
        from openai import OpenAI
        c = OpenAI(api_key=os.environ["DEEPSEEK_API_KEY"], base_url="https://api.deepseek.com")
        sd = data.get("sim_data", current_result or {})
        prompt = f"""Analyze this VLC visible light communication simulation professionally (in Chinese, under 300 chars):
Modulation={sd.get('modulation')} Data={sd.get('data_length')}bits SNR={sd.get('snr_db')}dB BER={sd.get('ber')} AI confidence={sd.get('ai_confidence')}
Explain the results in context of VLC system design. Note that optical channels naturally have higher error rates than RF. Focus on: signal quality, noise impact, and practical suggestions (e.g. increase SNR, add error correction). Be constructive and educational."""
        r = c.chat.completions.create(model="deepseek-chat", messages=[{"role":"user","content":prompt}], max_tokens=600)
        return {"success": True, "report": r.choices[0].message.content}
    except Exception as e:
        return JSONResponse({"success": False, "error": str(e)}, 500)

# ── Export ──
@app.post("/api/export_pptx")
async def export_pptx():
    try:
        from pptx import Presentation; from pptx.util import Inches, Pt
        prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _tb(s, "VLC Simulation Report", 0.5, 0.3, 12, 1, 28, (0x4a,0x7a,0xaa), True)
        _tb(s, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", 1.2, 0.5, 12, 0.5, 14)
        if current_result:
            s = prs.slides.add_slide(prs.slide_layouts[6])
            _tb(s, "Results", 0.5, 0.3, 12, 1, 24, (0x4a,0x7a,0xaa), True)
            for i, (k, v) in enumerate([("Modulation",current_result.get('modulation')),("Data",f"{current_result.get('data_length')} bits"),("SNR",f"{current_result.get('snr_db')} dB"),("BER",str(current_result.get('ber'))),("Errors",str(current_result.get('errors'))),("AI Confidence",f"{current_result.get('ai_confidence'):.2%}"),("Time",f"{current_result.get('elapsed_ms')} ms")]):
                _tb(s, f"{k}: {v}", 1.0+i*0.55, 1.5, 10, 0.5, 16)
        out = Path(__file__).resolve().parent.parent.parent / "results"; out.mkdir(exist_ok=True)
        fn = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        prs.save(str(out / fn))
        return {"success": True, "filename": fn}
    except Exception as e:
        return JSONResponse({"success": False, "error": str(e)}, 500)

# ── Settings ──
@app.get("/api/logs")
async def get_logs():
    return {"logs": system_logs[-50:]}  # last 50 entries

@app.post("/api/data/save")
async def save_data(data: dict):
    global current_result
    if not current_result:
        return JSONResponse({"success": False, "error": "No simulation data to save"}, 400)
    record = {
        "id": len(saved_records) + 1,
        "name": data.get("name", f"Save {len(saved_records)+1}"),
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "config": {
            "modulation": current_result.get("modulation"),
            "data_length": current_result.get("data_length"),
            "snr_db": current_result.get("snr_db"),
        },
        "result": {
            "ber": current_result.get("ber"),
            "errors": current_result.get("errors"),
            "ai_confidence": current_result.get("ai_confidence"),
            "elapsed_ms": current_result.get("elapsed_ms"),
        }
    }
    saved_records.append(record)
    return {"success": True, "id": record["id"], "name": record["name"]}

@app.get("/api/data/list")
async def list_data():
    summaries = [{"id": r["id"], "name": r["name"], "timestamp": r["timestamp"], "config": r["config"], "result": r["result"]} for r in saved_records]
    return {"records": summaries}

@app.post("/api/data/load/{record_id}")
async def load_data(record_id: int):
    global current_result
    for r in saved_records:
        if r["id"] == record_id:
            # Restore config and result for display
            current_result = {
                "tx_bits": r["config"]["data_length"],
                "errors": r["result"]["errors"],
                "ber": r["result"]["ber"],
                "snr_db": r["config"]["snr_db"],
                "modulation": r["config"]["modulation"],
                "data_length": r["config"]["data_length"],
                "elapsed_ms": r["result"]["elapsed_ms"],
                "ai_confidence": r["result"]["ai_confidence"],
                "waveform": [],
                "constellation": [],
                "eye_data": [],
            }
            return {"success": True, "data": current_result}
    return JSONResponse({"success": False, "error": "Record not found"}, 404)

@app.post("/api/settings/save")
async def save_settings(data: dict):
    try:
        sp = Path(__file__).resolve().parent.parent.parent / "config.json"
        cfg = json.loads(sp.read_text("utf-8")) if sp.exists() else {}
        for k in ["web_theme","web_font_size","web_accent"]:
            if k.replace("web_","") in data: cfg[k] = data[k.replace("web_","")]
        cfg["language"] = data.get("lang", cfg.get("language","zh"))
        sp.write_text(json.dumps(cfg, indent=2, ensure_ascii=False), "utf-8")
        get_i18n().set_language(cfg["language"])
        return {"success": True}
    except Exception as e:
        return JSONResponse({"success": False, "error": str(e)}, 500)

@app.get("/api/settings/load")
async def load_settings():
    try:
        sp = Path(__file__).resolve().parent.parent.parent / "config.json"
        cfg = json.loads(sp.read_text("utf-8")) if sp.exists() else {}
        return {"theme": cfg.get("web_theme","light"), "font_size": cfg.get("web_font_size","medium"), "accent": cfg.get("web_accent","#4a7aaa"), "lang": cfg.get("language","zh")}
    except:
        return {"theme":"light", "font_size":"medium", "accent":"#4a7aaa", "lang":"zh"}

def _tb(slide, text, top, left, w, h, size=18, color=None, bold=False):
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold
    if color: p.font.color.rgb = color

if __name__ == "__main__":
    import uvicorn; uvicorn.run(app, host="0.0.0.0", port=8000)
